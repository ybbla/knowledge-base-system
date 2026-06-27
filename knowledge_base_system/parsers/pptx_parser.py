"""PPTX 演示文稿解析器
使用 python-pptx .pptx 文件解析为统一ParsedElement Asset支持幻灯片标题、文本形状、列表、表格、图片和超链接资源的提取"""

import hashlib
import io
import re
import zipfile
from dataclasses import dataclass, field
from typing import Any

from app.core.models import (
    Asset,
    AssetData,
    AssetStatus,
    AssetType,
    Document,
    ElementType,
    ParsedElement,
    SourceLocation,
    compute_hash,
)
from parsers.base import DocumentParser, ParseResult
from parsers.utils import classify_link


@dataclass
class _AssetRecord:
    """内部资源记录，包含 Asset 和去重键。"""
    asset: Asset
    key: tuple[str, str]


@dataclass
class _ShapeRecord:
    """形状记录，包含排序所需的位置和索引。"""
    shape: Any
    index: int
    left: int
    top: int
    width: int
    height: int


@dataclass
class _PptxParseState:
    """PPTX 解析过程中的可变状态。"""
    doc_id: str
    doc_version: int
    elements: list[ParsedElement] = field(default_factory=list)
    assets: list[Asset] = field(default_factory=list)
    assets_by_key: dict[tuple[str, str], _AssetRecord] = field(default_factory=dict)
    seq: int = 0
    section_path: list[str] = field(default_factory=list)
    _counters: dict[str, int] = field(default_factory=dict)

    def next_seq(self) -> int:
        """生成递增的元素序号。"""
        self.seq += 1
        return self.seq

    _PH_MAP = {"image": "image", "video": "video", "image_link": "image",
               "video_link": "video", "document_link": "doc", "web_link": "web"}

    def next_ph(self, asset_type: str) -> str:
        """生成递增占位符 {{type:n}}。"""
        prefix = self._PH_MAP.get(asset_type, "res")
        self._counters[prefix] = self._counters.get(prefix, 0) + 1
        return f"{{{{{prefix}:{self._counters[prefix]}}}}}"


class PptxParser(DocumentParser):
    """将 PPTX 演示文稿解析为统一的 ParsedElement 和 Asset。

    按幻灯片顺序处理标题、列表、表格、图片和链接资源。
    """

    SUPPORTED_TYPES = {"pptx"}
    TITLE_PLACEHOLDER_TYPES = {"TITLE", "CENTER_TITLE"}
    BODY_PLACEHOLDER_TYPES = {"BODY", "OBJECT", "CONTENT"}
    UNSUPPORTED_SHAPE_TYPES = {
        "EMBEDDED_OLE_OBJECT",
        "GROUP",
        "LINKED_OLE_OBJECT",
    }

    def supports(self, source_type: str) -> bool:
        return source_type.lower() in self.SUPPORTED_TYPES

    def parse(self, doc: Document, content: bytes | str) -> ParseResult:
        """将 PPTX 文档解析为结构化元素和资源列表。"""
        if isinstance(content, str):
            content = content.encode("utf-8")
        if not content:
            raise ValueError("PPTX 解析失败：文档内容为空")

        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("PPTX 解析失败：缺python-pptx 依赖") from exc

        try:
            presentation = Presentation(io.BytesIO(content))
        except (OSError, KeyError, ValueError, zipfile.BadZipFile) as exc:
            raise ValueError(f"PPTX 解析失败：{exc}") from exc
        except Exception as exc:
            raise ValueError(f"PPTX 解析失败：{exc}") from exc

        if len(presentation.slides) == 0:
            raise ValueError("PPTX 解析失败：未包含可解析幻灯片")

        state = _PptxParseState(doc.doc_id, doc.version)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            self._process_slide(slide, slide_index, state, doc)

        if not state.elements:
            raise ValueError("PPTX 解析失败：未提取到有效内容")

        doc.source_hash = compute_hash(content)
        return ParseResult(doc=doc, elements=state.elements, assets=state.assets)

    def _process_slide(
        self,
        slide: Any,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> None:
        """处理单张幻灯片：识别标题后按顺序处理其余形状。"""
        shape_records = self._shape_records(slide)
        title_record = self._find_title_shape(shape_records)
        title_text = (
            self._shape_text(title_record.shape)
            if title_record is not None
            else ""
        ) or f"幻灯{slide_index}"

        state.section_path = [title_text]
        self._append_element(
            state,
            ParsedElement(
                doc_id=state.doc_id,
                doc_version=state.doc_version,
                sequence_order=state.next_seq(),
                element_type=ElementType.title,
                text=title_text,
                source_location=SourceLocation(section_path=list(state.section_path)),
                metadata={
                    **self._slide_metadata(slide_index),
                    **(
                        self._shape_metadata(title_record)
                        if title_record is not None
                        else {}
                    ),
                    "heading_level": 1,
                    "fallback_title": title_record is None,
                },
            ),
        )

        for record in shape_records:
            if title_record is not None and record.shape is title_record.shape:
                continue
            if self._shape_has_table(record.shape):
                self._add_table(record, slide_index, state, doc, slide)
            elif self._is_picture(record.shape):
                self._add_image(record, slide_index, state, doc)
            elif self._shape_has_text(record.shape):
                self._add_text_shape(record, slide_index, state, doc)
            elif self._is_chart_shape(record.shape):
                self._add_chart(record, slide_index, state, doc)
            elif self._is_video_shape(record.shape):
                self._add_video(record, slide_index, state, doc, slide)
            elif self._is_unsupported_shape(record.shape):
                self._add_unknown(record, slide_index, state)

    def _shape_records(self, slide: Any) -> list[_ShapeRecord]:
        """按顶部、左侧和原始序号收集并排序形状。"""
        records: list[_ShapeRecord] = []
        for index, shape in enumerate(slide.shapes):
            records.append(
                _ShapeRecord(
                    shape=shape,
                    index=index,
                    left=self._coord(shape, "left"),
                    top=self._coord(shape, "top"),
                    width=self._coord(shape, "width"),
                    height=self._coord(shape, "height"),
                )
            )
        return sorted(records, key=lambda item: (item.top, item.left, item.index))

    def _find_title_shape(self, records: list[_ShapeRecord]) -> _ShapeRecord | None:
        """定位幻灯片的标题形状
        优先匹配 TITLE/CENTER_TITLE 占位符类型，其次选首个有文本的形状        """
        for record in records:
            if not self._shape_has_text(record.shape):
                continue
            if self._placeholder_type_name(record.shape) in self.TITLE_PLACEHOLDER_TYPES:
                if self._shape_text(record.shape):
                    return record

        for record in records:
            if self._shape_has_text(record.shape) and self._shape_text(record.shape):
                return record
        return None

    def _add_text_shape(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> None:
        """处理文本形状：遍历 run 级别，仅处理超链接，锚定文字用占位符替换。"""
        if not self._shape_has_text(record.shape):
            return

        # 按段落处理：遍历 run，run 超链接 → 占位符
        paragraphs_data: list[dict[str, Any]] = []
        for paragraph in record.shape.text_frame.paragraphs:
            text, ads = self._render_paragraph_runs(paragraph, record, slide_index, state, doc)
            if text:
                paragraphs_data.append({
                    "text": text,
                    "level": getattr(paragraph, "level", 0),
                    "asset_data": ads,
                })

        if not paragraphs_data:
            return

        # 形状级超链接 → 挂到第一个段落
        try:
            shape_url = record.shape.click_action.hyperlink.address
        except Exception:
            shape_url = None
        if shape_url:
            asset = self._asset_for_url(
                shape_url, self._asset_type_for_url(shape_url), state, doc,
                {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                 "source": "pptx_hyperlink"},
                display_text="",
            )
            ph = state.next_ph(asset.asset_type.value)
            paragraphs_data[0]["asset_data"].append(
                AssetData(placeholder=ph, asset_id=asset.asset_id))
            paragraphs_data[0]["text"] += " " + ph

        text = "\n".join(pd["text"] for pd in paragraphs_data)
        all_ads: list[AssetData] = []
        for pd in paragraphs_data:
            all_ads.extend(pd["asset_data"])
        element_type = (
            ElementType.list if self._is_list_shape(record.shape, paragraphs_data)
            else ElementType.paragraph
        )
        self._append_element(state, ParsedElement(
            doc_id=state.doc_id,
            doc_version=state.doc_version,
            sequence_order=state.next_seq(),
            element_type=element_type,
            text=text,
            asset_data=all_ads,
            source_location=SourceLocation(section_path=list(state.section_path)),
            metadata={
                **self._slide_metadata(slide_index),
                **self._shape_metadata(record),
            },
        ))

    def _render_paragraph_runs(
        self,
        paragraph: Any,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> tuple[str, list[AssetData]]:
        """处理单段落的 run 序列：仅处理 run 级超链接，不提取纯文本 URL。

        与 DOCX 对齐：仅真正的超链接对象创建 Asset，锚定文字被占位符替换。
        """
        text_parts: list[str] = []
        asset_data_list: list[AssetData] = []

        for run in paragraph.runs:
            run_text = self._normalize_text(run.text)
            if not run_text:
                continue

            try:
                address = run.hyperlink.address
            except Exception:
                address = None

            if address:
                asset = self._asset_for_url(
                    address, self._asset_type_for_url(address), state, doc,
                    {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                     "source": "pptx_hyperlink"},
                    display_text=run_text,
                )
                ph = state.next_ph(asset.asset_type.value)
                asset_data_list.append(AssetData(placeholder=ph, asset_id=asset.asset_id))
                text_parts.append(ph)
            else:
                text_parts.append(run_text)

        return "".join(text_parts), asset_data_list

    def _render_cell_text(
        self,
        cell: Any,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
        slide: Any = None,
    ) -> tuple[str, list[AssetData]]:
        """处理表格单元格文本：仅处理 run 级超链接，不提取纯文本 URL。

        text_frame 异常时 fallback 到 raw XML，确保超链接不丢失。
        """
        text_parts: list[str] = []
        asset_data_list: list[AssetData] = []

        try:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run_text = self._normalize_text(run.text)
                    if not run_text:
                        continue
                    try:
                        addr = run.hyperlink.address
                    except Exception:
                        addr = None
                    if addr:
                        asset = self._asset_for_url(
                            addr, self._asset_type_for_url(addr), state, doc,
                            {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                             "source": "pptx_table_cell"},
                            display_text=run_text,
                        )
                        ph = state.next_ph(asset.asset_type.value)
                        asset_data_list.append(AssetData(placeholder=ph, asset_id=asset.asset_id))
                        text_parts.append(ph)
                    else:
                        text_parts.append(run_text)
        except Exception:
            if slide is not None:
                return self._render_cell_xml_fallback(cell, slide, record, slide_index, state, doc)
            return self._normalize_text(cell.text), []

        return "".join(p for p in text_parts if p), asset_data_list

    def _render_cell_xml_fallback(
        self,
        cell: Any,
        slide: Any,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> tuple[str, list[AssetData]]:
        """text_frame 不可用时从 raw XML 解析单元格 run 级超链接。"""
        A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        text_parts: list[str] = []
        asset_data_list: list[AssetData] = []
        rels = slide.part.rels

        for p_el in cell._tc.iter(f"{{{A_NS}}}p"):
            for r_el in p_el.iter(f"{{{A_NS}}}r"):
                t_el = r_el.find(f"{{{A_NS}}}t")
                if t_el is None or not t_el.text:
                    continue
                run_text = self._normalize_text(t_el.text)
                if not run_text:
                    continue

                rPr = r_el.find(f"{{{A_NS}}}rPr")
                hlink = rPr.find(f"{{{A_NS}}}hlinkClick") if rPr is not None else None
                if hlink is not None:
                    rId = hlink.get(f"{{{R_NS}}}id")
                    if rId:
                        try:
                            url = rels[rId].target_ref
                        except (KeyError, AttributeError):
                            url = None
                        if url:
                            asset = self._asset_for_url(
                                url, self._asset_type_for_url(url), state, doc,
                                {**self._slide_metadata(slide_index),
                                 **self._shape_metadata(record),
                                 "source": "pptx_table_cell"},
                                display_text=run_text,
                            )
                            ph = state.next_ph(asset.asset_type.value)
                            asset_data_list.append(
                                AssetData(placeholder=ph, asset_id=asset.asset_id))
                            text_parts.append(ph)
                            continue

                text_parts.append(run_text)

        return "".join(text_parts), asset_data_list

    def _add_table(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
        slide: Any = None,
    ) -> None:
        """处理表格形状：仅处理 run 超链接，锚定文字用占位符替换，与 DOCX 对齐。"""
        table = record.shape.table
        rows: list[list[dict[str, Any]]] = []
        all_asset_data: list[AssetData] = []

        for row_index, row in enumerate(table.rows, start=1):
            cells = []
            for col_index, cell in enumerate(row.cells, start=1):
                text, cell_ads = self._render_cell_text(cell, record, slide_index, state, doc, slide)
                all_asset_data.extend(cell_ads)
                cells.append({
                    "text": text,
                    "metadata": {
                        "row": row_index,
                        "column": col_index,
                        "slide_index": slide_index,
                        "slide_number": slide_index,
                    },
                })
            rows.append(cells)

        if not rows:
            return

        # ── 形状级超链接 ──
        try:
            addr = record.shape.click_action.hyperlink.address
        except Exception:
            addr = None
        if addr:
            asset = self._asset_for_url(
                addr, self._asset_type_for_url(addr), state, doc,
                {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                 "source": "pptx_table_hyperlink"},
                display_text="",
            )
            ph = state.next_ph(asset.asset_type.value)
            all_asset_data.append(AssetData(placeholder=ph, asset_id=asset.asset_id))
            # 占位符追加到第一个单元格文本
            rows[0][0]["text"] += " " + ph

        headers = [{"text": cell["text"]} for cell in rows[0]]
        data_rows = [{"cells": row} for row in rows[1:]]
        structured: dict[str, Any] = {
            "table": {
                "caption": "",
                "headers": headers,
                "rows": data_rows,
                "metadata": {
                    **self._slide_metadata(slide_index),
                    **self._shape_metadata(record),
                    "row_count": len(rows),
                    "column_count": len(rows[0]) if rows else 0,
                },
            }
        }

        text_parts = []
        if headers:
            text_parts.append(" | ".join(h["text"] for h in headers))
        for row in data_rows:
            text_parts.append(" | ".join(cell["text"] for cell in row["cells"]))

        metadata = {
            **self._slide_metadata(slide_index),
            **self._shape_metadata(record),
        }

        self._append_element(
            state,
            ParsedElement(
                doc_id=state.doc_id,
                doc_version=state.doc_version,
                sequence_order=state.next_seq(),
                element_type=ElementType.table,
                text="\n".join(part for part in text_parts if part.strip()),
                structured_data=structured,
                asset_data=all_asset_data,
                source_location=SourceLocation(section_path=list(state.section_path)),
                metadata=metadata,
            ),
        )

    def _add_image(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> None:
        """处理图片形状，创建图片 Asset 和对应元素。"""
        image = record.shape.image
        data = image.blob
        content_hash = f"sha256:{hashlib.sha256(data).hexdigest()}"
        filename = image.filename or f"image-{slide_index}-{record.index + 1}.{image.ext}"
        key = ("image", content_hash)
        record_asset = state.assets_by_key.get(key)
        if record_asset is None:
            asset = Asset(
                doc_id=doc.doc_id,
                asset_type=AssetType.image,
                original_uri=f"pptx://{doc.doc_id}/slide/{slide_index}/media/{filename}",
                storage_uri=None,
                content_hash=content_hash,
                status=AssetStatus.ready,
                extracted_text=None,
                metadata={
                    **self._slide_metadata(slide_index),
                    **self._shape_metadata(record),
                    "mime_type": image.content_type,
                    "file_name": filename,
                    "source": "pptx_image",
                },
            )
            object.__setattr__(asset, "_data", data)
            state.assets.append(asset)
            state.assets_by_key[key] = _AssetRecord(asset=asset, key=key)
        else:
            asset = record_asset.asset

        # 图片自身的 AssetData
        img_ph = state.next_ph("image")
        text = img_ph
        all_ads: list[AssetData] = [
            AssetData(placeholder=img_ph, asset_id=asset.asset_id)
        ]

        # 形状级超链接：Asset → 占位符追加到文本
        try:
            shape_url = record.shape.click_action.hyperlink.address
        except Exception:
            shape_url = None
        if shape_url:
            link_asset = self._asset_for_url(
                shape_url, self._asset_type_for_url(shape_url), state, doc,
                {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                 "source": "pptx_hyperlink"},
                display_text="",
            )
            link_ph = state.next_ph(link_asset.asset_type.value)
            all_ads.append(AssetData(placeholder=link_ph, asset_id=link_asset.asset_id))
            text += " " + link_ph

        # 运行级超链接
        if self._shape_has_text(record.shape):
            for paragraph in record.shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    try:
                        addr = run.hyperlink.address
                    except Exception:
                        addr = None
                    if addr:
                        run_text = self._normalize_text(run.text)
                        run_asset = self._asset_for_url(
                            addr, self._asset_type_for_url(addr), state, doc,
                            {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                             "source": "pptx_hyperlink"},
                            display_text=run_text,
                        )
                        run_ph = state.next_ph(run_asset.asset_type.value)
                        all_ads.append(AssetData(placeholder=run_ph, asset_id=run_asset.asset_id))
                        text += " " + run_ph

        self._append_element(
            state,
            ParsedElement(
                doc_id=state.doc_id,
                doc_version=state.doc_version,
                sequence_order=state.next_seq(),
                element_type=ElementType.paragraph,
                text=text,
                asset_data=all_ads,
                source_location=SourceLocation(section_path=list(state.section_path)),
                metadata={
                    **self._slide_metadata(slide_index),
                    **self._shape_metadata(record),
                    "file_name": filename,
                },
            ),
        )

    def _add_video(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
        slide: Any,
    ) -> None:
        """处理视频形状：从 raw XML 提取外部视频链接，创建 video_link Asset。"""
        A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        url = None
        el = record.shape._element
        video_file = el.find(f".//{{{A_NS}}}videoFile")
        if video_file is not None:
            rId = video_file.get(f"{{{R_NS}}}link")
            if rId:
                try:
                    url = slide.part.rels[rId].target_ref
                except (KeyError, AttributeError):
                    pass

        if not url:
            self._add_unknown(record, slide_index, state)
            return

        asset = self._asset_for_url(
            url, AssetType.video_link, state, doc,
            {**self._slide_metadata(slide_index), **self._shape_metadata(record),
             "source": "pptx_video"},
        )
        ph = state.next_ph(asset.asset_type.value)
        text = ph
        all_ads = [AssetData(placeholder=ph, asset_id=asset.asset_id)]

        # 形状级超链接
        try:
            shape_url = record.shape.click_action.hyperlink.address
        except Exception:
            shape_url = None
        if shape_url:
            link_asset = self._asset_for_url(
                shape_url, self._asset_type_for_url(shape_url), state, doc,
                {**self._slide_metadata(slide_index), **self._shape_metadata(record),
                 "source": "pptx_hyperlink"},
                display_text="",
            )
            link_ph = state.next_ph(link_asset.asset_type.value)
            all_ads.append(AssetData(placeholder=link_ph, asset_id=link_asset.asset_id))
            text += " " + link_ph

        self._append_element(state, ParsedElement(
            doc_id=state.doc_id,
            doc_version=state.doc_version,
            sequence_order=state.next_seq(),
            element_type=ElementType.paragraph,
            text=text,
            asset_data=all_ads,
            source_location=SourceLocation(section_path=list(state.section_path)),
            metadata={
                **self._slide_metadata(slide_index),
                **self._shape_metadata(record),
            },
        ))

    @staticmethod
    def _is_video_shape(shape: Any) -> bool:
        """判断形状是否为视频（MEDIA 类型）。"""
        shape_type = getattr(shape, "shape_type", "")
        type_name = getattr(shape_type, "name", str(shape_type)).upper()
        return type_name == "MEDIA"

    def _add_chart(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
        doc: Document,
    ) -> None:
        """处理图表形状：提取标题和表格数据，生成为 title + table 元素。"""
        chart = record.shape.chart

        # ── 图表标题 ──
        chart_title = ""
        try:
            if chart.chart_title is not None and chart.chart_title.has_text_frame:
                chart_title = self._normalize_text(chart.chart_title.text_frame.text)
        except Exception:
            pass

        # ── 提取分类和数据系列 ──
        categories: list[str] = []
        try:
            cats = chart.plots[0].categories
            if cats is not None:
                categories = [str(c) for c in cats]
        except Exception:
            try:
                cats = getattr(chart.category_axis, "categories", None)
                if cats is not None:
                    categories = [str(c) for c in cats]
            except Exception:
                pass

        series_names: list[str] = []
        series_values: list[list] = []
        try:
            for s in chart.series:
                name = ""
                try:
                    if s.name is not None:
                        name = str(s.name)
                except Exception:
                    pass
                series_names.append(name)
                try:
                    series_values.append(list(s.values))
                except Exception:
                    series_values.append([])
        except Exception:
            pass

        if not categories and not series_names:
            self._add_unknown(record, slide_index, state)
            return

        # ── 构建表格 ──
        # headers: [类别, 系列1, 系列2, ...]
        headers_data = [{"text": chart_title or "类别"}]
        for name in series_names:
            headers_data.append({"text": name})

        rows_data = []
        for i, cat in enumerate(categories):
            cells = [{"text": cat}]
            for vals in series_values:
                cells.append({"text": str(vals[i]) if i < len(vals) else ""})
            rows_data.append({"cells": cells})

        all_asset_data: list[AssetData] = []

        text_parts = [" | ".join(h["text"] for h in headers_data)]
        for row in rows_data:
            text_parts.append(" | ".join(c["text"] for c in row["cells"]))

        self._append_element(state, ParsedElement(
            doc_id=state.doc_id,
            doc_version=state.doc_version,
            sequence_order=state.next_seq(),
            element_type=ElementType.table,
            text="\n".join(text_parts),
            structured_data={
                "table": {
                    "caption": "",
                    "headers": headers_data,
                    "rows": rows_data,
                    "metadata": {
                        **self._slide_metadata(slide_index),
                        **self._shape_metadata(record),
                        "source": "pptx_chart",
                    },
                }
            },
            asset_data=all_asset_data,
            source_location=SourceLocation(section_path=list(state.section_path)),
            metadata={
                **self._slide_metadata(slide_index),
                **self._shape_metadata(record),
            },
        ))

    @staticmethod
    def _is_chart_shape(shape: Any) -> bool:
        """判断形状是否为图表。"""
        return bool(getattr(shape, "has_chart", False))

    def _add_unknown(
        self,
        record: _ShapeRecord,
        slide_index: int,
        state: _PptxParseState,
    ) -> None:
        """为图表、OLE 等不支持的形状添加占位元素。"""
        self._append_element(
            state,
            ParsedElement(
                doc_id=state.doc_id,
                doc_version=state.doc_version,
                sequence_order=state.next_seq(),
                element_type=ElementType.unknown,
                text=f"[不支持的 PPTX 对象: {self._shape_type_name(record.shape)}]",
                source_location=SourceLocation(section_path=list(state.section_path)),
                metadata={
                    **self._slide_metadata(slide_index),
                    **self._shape_metadata(record),
                },
            ),
        )

    def _asset_for_url(
        self,
        url: str,
        asset_type: AssetType,
        state: _PptxParseState,
        doc: Document,
        metadata: dict[str, Any],
        display_text: str = "",
    ) -> Asset:
        """创建或复用 Asset，按 URL 和资源类型去重。

        首次创建时设置 display_text（超链接锚定文字），后续复用不覆盖。
        """
        key = (asset_type.value, url)
        existing = state.assets_by_key.get(key)
        if existing is not None:
            return existing.asset

        asset = Asset(
            doc_id=doc.doc_id,
            asset_type=asset_type,
            original_uri=url,
            display_text=display_text,
            storage_uri=None,
            status=AssetStatus.ready,
            extracted_text=None,
            metadata={**metadata},
        )
        state.assets.append(asset)
        state.assets_by_key[key] = _AssetRecord(asset=asset, key=key)
        return asset

    @staticmethod
    def _append_element(state: _PptxParseState, element: ParsedElement) -> None:
        """添加元素，并通过 asset_id 回填 Asset.element_id。"""
        linked = {ad.asset_id for ad in element.asset_data}
        for record in state.assets_by_key.values():
            if record.asset.asset_id in linked and not record.asset.element_id:
                record.asset.element_id = element.element_id
        state.elements.append(element)

    def _is_list_shape(self, shape: Any, paragraphs: list[dict[str, Any]]) -> bool:
        """判断形状内容是否为列表格式
        依据：多段落 + 有缩进层级，或占位符类型BODY        """
        if len(paragraphs) <= 1:
            return False
        if any(item["level"] > 0 for item in paragraphs):
            return True
        return self._placeholder_type_name(shape) in self.BODY_PLACEHOLDER_TYPES

    @staticmethod
    def _shape_has_text(shape: Any) -> bool:
        """判断形状是否包含文本框。"""
        return bool(getattr(shape, "has_text_frame", False))

    @staticmethod
    def _shape_has_table(shape: Any) -> bool:
        """判断形状是否包含表格。"""
        return bool(getattr(shape, "has_table", False))

    def _shape_text(self, shape: Any) -> str:
        """获取形状的纯文本内容并归一化空白。"""
        if not self._shape_has_text(shape):
            return ""
        return self._normalize_text(shape.text)

    @staticmethod
    def _normalize_text(text: str) -> str:
        """将连续空白归一化为单个空格并去除首尾空白。"""
        return re.sub(r"[ \t\r\f\v]+", " ", text or "").strip()

    def _is_picture(self, shape: Any) -> bool:
        """判断形状是否为图片。"""
        if hasattr(shape, "image"):
            return True
        return self._shape_type_name(shape) == "PICTURE"

    def _is_unsupported_shape(self, shape: Any) -> bool:
        """判断形状是否为图表、OLE 等不支持的类型。"""
        return self._shape_type_name(shape) in self.UNSUPPORTED_SHAPE_TYPES

    @staticmethod
    def _shape_type_name(shape: Any) -> str:
        """获取形状类型的大写名称。"""
        value = getattr(shape, "shape_type", "")
        return getattr(value, "name", str(value)).upper()

    @staticmethod
    def _placeholder_type_name(shape: Any) -> str:
        """获取占位符类型的大写名称。"""
        if not getattr(shape, "is_placeholder", False):
            return ""
        try:
            value = shape.placeholder_format.type
        except Exception:
            return ""
        return getattr(value, "name", str(value)).upper()

    @staticmethod
    def _coord(shape: Any, attr: str) -> int:
        """安全获取形状的 EMU 坐标整数值。"""
        value = getattr(shape, attr, 0)
        try:
            return int(value)
        except (TypeError, ValueError):
            return 0

    @staticmethod
    def _slide_metadata(slide_index: int) -> dict[str, Any]:
        """生成幻灯片元数据字典。"""
        return {"slide_index": slide_index, "slide_number": slide_index}

    def _shape_metadata(self, record: _ShapeRecord) -> dict[str, Any]:
        """生成包含 ID、名称、类型和位置的形状元数据。"""
        shape = record.shape
        return {
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": getattr(shape, "name", ""),
            "shape_type": self._shape_type_name(shape),
            "placeholder_type": self._placeholder_type_name(shape),
            "shape_index": record.index,
            "left": record.left,
            "top": record.top,
            "width": record.width,
            "height": record.height,
        }

    def _asset_type_for_url(self, url: str) -> AssetType | None:
        """根据 URL 模式判断资源类型（使用classify_link 分类）。

        与 DOCX 解析器对齐：所有 URL 类型均创建 Asset。
        """
        kind = classify_link(url)
        return {
            "image": AssetType.image_link,
            "video": AssetType.video_link,
            "audio": AssetType.video_link,
            "document": AssetType.document_link,
            "url": AssetType.web_link,
        }.get(kind)  # 都返回 AssetType，不再有 None

