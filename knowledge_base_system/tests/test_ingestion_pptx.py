import io

from pptx import Presentation
from pptx.util import Inches

from app.core.models import Document, ElementType, KnowledgeChunk
from ingestion.pipeline import IngestionPipeline
import ingestion.pipeline as ingestion_module
from parsers.pptx_parser import PptxParser
from parsers.registry import ParserRegistry


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "演示文稿"
    slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(0.5)).text = "PPTX 文档解析后进入语义抽取流程。"
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class _FakeExtractor:
    def __init__(self) -> None:
        self.seen_elements = []

    def extract(self, elements, assets, doc_id, category):
        self.seen_elements = list(elements)
        source = elements[0]
        return [
            KnowledgeChunk(
                doc_id=source.doc_id,
                doc_version=source.doc_version,
                title="PPTX 演示文稿",
                content="PPTX 文档解析后进入语义抽取流程。",
                category=category,
            )
        ]


class _RecordingIndex:
    def __init__(self) -> None:
        self.batches = []

    def add_batch(self, items):
        self.batches.append(items)

    def add(self, *_args, **_kwargs):
        raise AssertionError("入库管线应使用批量索引写入")

    def delete(self, _chunk_id):
        return None

    def search(self, *_args, **_kwargs):
        return []


class _RecordingChunkStore:
    def __init__(self) -> None:
        self.chunks = {}

    def put(self, chunk):
        self.chunks[chunk.chunk_id] = chunk


class _RecordingAssetStore:
    def __init__(self) -> None:
        self.assets = {}

    def put(self, asset):
        self.assets[asset.asset_id] = asset

    def get(self, asset_id):
        return self.assets.get(asset_id)

    def delete(self, asset_id):
        self.assets.pop(asset_id, None)


class _FakeEmbedder:
    def __init__(self) -> None:
        self.calls = []

    def embed_text(self, texts):
        self.calls.append(list(texts))
        return [[1.0, 0.0] for _ in texts]


def test_ingestion_pipeline_dispatches_pptx_parser(monkeypatch):
    embedder = _FakeEmbedder()
    monkeypatch.setattr(ingestion_module, "embedding_client", embedder)

    registry = ParserRegistry()
    registry.register(PptxParser())
    extractor = _FakeExtractor()
    vector_index = _RecordingIndex()
    bm25_index = _RecordingIndex()
    chunk_store = _RecordingChunkStore()
    asset_store = _RecordingAssetStore()
    pipeline = IngestionPipeline(
        parser_registry=registry,
        extractor=extractor,
        vector_index=vector_index,
        bm25_index=bm25_index,
        asset_store=asset_store,
        chunk_store=chunk_store,
    )
    doc = Document(
        title="演示文稿",
        source_type="pptx",
        source_uri="memory://slides.pptx",
        category="验收",
        metadata={"raw_content": _pptx_bytes()},
    )
    doc = pipeline.ingest(doc)

    assert doc.status.value == "active", doc.error_message
    assert [el.element_type for el in extractor.seen_elements] == [
        ElementType.title,
        ElementType.paragraph,
    ]
    assert len(chunk_store.chunks) == 1
    assert len(embedder.calls) == 1
    assert len(vector_index.batches) == 1
    assert len(bm25_index.batches) == 1


def test_ingestion_pipeline_marks_invalid_pptx_failed(monkeypatch):
    embedder = _FakeEmbedder()
    monkeypatch.setattr(ingestion_module, "embedding_client", embedder)

    pipeline = IngestionPipeline(
        parser_registry=ParserRegistry(),
        extractor=_FakeExtractor(),
        vector_index=_RecordingIndex(),
        bm25_index=_RecordingIndex(),
        asset_store=_RecordingAssetStore(),
        chunk_store=_RecordingChunkStore(),
    )
    pipeline._parser_registry.register(PptxParser())
    doc = Document(
        title="坏文件",
        source_type="pptx",
        source_uri="memory://bad.pptx",
        metadata={"raw_content": b"not a pptx"},
    )
    doc = pipeline.ingest(doc)

    assert doc.status.value == "failed"
    assert "PPTX 解析失败" in doc.error_message
