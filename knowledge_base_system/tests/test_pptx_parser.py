import base64
import io
import os
import tempfile

import pytest
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app.core.models import AssetType, Document, ElementType
from parsers.pptx_parser import PptxParser


def _pptx_bytes(presentation: Presentation) -> bytes:
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _doc(raw: bytes, source_uri: str = "memory://slides.pptx") -> Document:
    return Document(
        title="PPTX",
        source_type="pptx",
        source_uri=source_uri,
    )


def _parse(parser: PptxParser, raw: bytes, source_uri: str = "memory://slides.pptx"):
    """辅助：同时传入 doc 和 content 调用 parser.parse。"""
    doc = _doc(raw, source_uri)
    return parser.parse(doc, raw)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


class TestPptxParser:
    def setup_method(self):
        self.parser = PptxParser()

    def test_supported_types(self):
        assert self.parser.supports("pptx")
        assert self.parser.supports("PPTX")
        assert not self.parser.supports("ppt")

    def test_parse_title_paragraph_list_and_hash(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5))
        title.text = "上传流程"
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.5))
        body.text_frame.text = "准备文件"
        item = body.text_frame.add_paragraph()
        item.text = "上传文件"
        item.level = 1

        result = _parse(self.parser, _pptx_bytes(prs))

        assert result.doc.source_hash.startswith("sha256:")
        assert [el.element_type for el in result.elements] == [
            ElementType.title,
            ElementType.list,
            ElementType.paragraph,
            ElementType.paragraph,
        ]
        assert result.elements[0].text == "上传流程"
        assert result.elements[0].source_location.section_path == ["上传流程"]
        list_el = result.elements[1]
        assert all(el.parent_element_id == list_el.element_id for el in result.elements[2:])
        assert [el.text for el in result.elements[2:]] == ["准备文件", "上传文件"]
        assert result.elements[2].metadata["slide_index"] == 1
        assert result.elements[3].metadata["level"] == 1

    def test_parse_table(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "状态表"
        table_shape = slide.shapes.add_table(
            3,
            2,
            Inches(0.5),
            Inches(1.0),
            Inches(6),
            Inches(1.5),
        )
        table = table_shape.table
        table.cell(0, 0).text = "状态"
        table.cell(0, 1).text = "说明"
        table.cell(1, 0).text = "处理中"
        table.cell(1, 1).text = "系统正在解析文档"
        table.cell(2, 0).text = "成功"
        table.cell(2, 1).text = "文档已经进入知识库"

        result = _parse(self.parser, _pptx_bytes(prs))
        table_el = next(el for el in result.elements if el.element_type == ElementType.table)
        table_data = table_el.structured_data["table"]

        assert table_data["headers"] == ["状态", "说明"]
        assert table_data["rows"][0]["cells"][0]["text"] == "处理中"
        assert table_data["rows"][0]["cells"][0]["metadata"]["row"] == 2
        assert table_data["metadata"]["slide_index"] == 1
        assert table_el.source_location.section_path == ["状态表"]

    def test_extract_image_asset(self):
        png_bytes = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
        )
        image_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        try:
            image_file.write(png_bytes)
            image_file.close()
            prs = Presentation()
            slide = _blank_slide(prs)
            slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "截图"
            slide.shapes.add_picture(image_file.name, Inches(0.5), Inches(1.0))
            result = _parse(self.parser, _pptx_bytes(prs))
        finally:
            os.unlink(image_file.name)

        assert len(result.assets) == 1
        assert result.assets[0].asset_type == AssetType.image
        assert result.assets[0].content_hash.startswith("sha256:")
        assert getattr(result.assets[0], "_data") == png_bytes
        image_el = next(el for el in result.elements if el.element_type == ElementType.paragraph)
        assert image_el.asset_data[0].url == result.assets[0].original_uri
        assert result.assets[0].element_id == image_el.element_id

    def test_video_audio_attachment_links_and_dedup(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "链接"
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.0))
        textbox.text = (
            "演示视频 https://example.com/demo.mp4 "
            "音频 https://example.com/audio.mp3 "
            "附件 https://example.com/manual.pdf "
            "重复 https://example.com/demo.mp4"
        )

        result = _parse(self.parser, _pptx_bytes(prs))
        assets = {asset.original_uri: asset for asset in result.assets}

        assert assets["https://example.com/demo.mp4"].asset_type == AssetType.video_link
        assert assets["https://example.com/audio.mp3"].asset_type == AssetType.video_link
        assert assets["https://example.com/manual.pdf"].asset_type == AssetType.document_link
        assert len(result.assets) == 3
        paragraph = next(
            el
            for el in result.elements
            if el.element_type == ElementType.paragraph and "演示视频" in el.text
        )
        assert len(paragraph.asset_data) == 3
        assert all(asset.element_id == paragraph.element_id for asset in result.assets)

    def test_unsupported_chart_becomes_unknown_element(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "图表页"
        chart_data = ChartData()
        chart_data.categories = ["一月", "二月"]
        chart_data.add_series("数量", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5),
            Inches(1.0),
            Inches(6),
            Inches(3),
            chart_data,
        )

        result = _parse(self.parser, _pptx_bytes(prs))
        unknown = next(el for el in result.elements if el.element_type == ElementType.unknown)

        assert "不支持的 PPTX 对象" in unknown.text
        assert unknown.metadata["shape_type"] == "CHART"

    def test_blank_slide_gets_fallback_title(self):
        prs = Presentation()
        _blank_slide(prs)

        result = _parse(self.parser, _pptx_bytes(prs))

        assert len(result.elements) == 1
        assert result.elements[0].element_type == ElementType.title
        assert result.elements[0].text == "幻灯片 1"
        assert result.elements[0].metadata["fallback_title"] is True

    def test_read_from_file_uri(self, tmp_path):
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "文件来源"
        path = tmp_path / "sample.pptx"
        raw = _pptx_bytes(prs)
        path.write_bytes(raw)

        doc = Document(title="File", source_type="pptx", source_uri=f"file://{path}")
        result = self.parser.parse(doc, raw)

        assert result.doc.source_hash.startswith("sha256:")
        assert result.elements[0].text == "文件来源"

    def test_invalid_pptx_raises_clear_error(self):
        with pytest.raises(ValueError, match="PPTX 解析失败"):
            _parse(self.parser, b"not a pptx")

    def test_empty_presentation_raises_clear_error(self):
        prs = Presentation()
        with pytest.raises(ValueError, match="PPTX 解析失败"):
            _parse(self.parser, _pptx_bytes(prs))

    # ── 新增：超链接文字保留和 structured_data.links 测试 ──────────────

    def test_hyperlink_run_text_preserved_in_links(self):
        """验证文本运行中超链接的文字保留和 structured_data.links 输出。"""
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "链接测试"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.0))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "普通文字"
        run2 = p.add_run()
        run2.text = "点击查看文档"
        run2.hyperlink.address = "https://example.com/doc.pdf"

        result = _parse(self.parser, _pptx_bytes(prs))
        paragraph = next(
            el for el in result.elements
            if el.element_type == ElementType.paragraph and "普通文字" in el.text
        )
        assert "点击查看文档" in paragraph.text
        assert paragraph.structured_data is not None
        links = paragraph.structured_data["links"]
        assert len(links) == 1
        assert links[0]["text"] == "点击查看文档"
        assert links[0]["url"] == "https://example.com/doc.pdf"
        assert links[0]["link_type"] == "document"
        # 验证 Asset 被创建
        assert len(result.assets) == 1
        assert result.assets[0].original_uri == "https://example.com/doc.pdf"
        assert result.assets[0].asset_type == AssetType.document_link

    def test_shape_level_hyperlink_preserved_in_links(self):
        """验证形状级超链接（click_action.hyperlink）的文字和链接记录。"""
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "视频页"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(0.8))
        tb.text_frame.paragraphs[0].text = "观看演示"
        tb.click_action.hyperlink.address = "https://example.com/demo.mp4"

        result = _parse(self.parser, _pptx_bytes(prs))
        paragraph = next(
            el for el in result.elements
            if el.element_type == ElementType.paragraph and "观看演示" in el.text
        )
        assert paragraph.structured_data is not None
        links = paragraph.structured_data["links"]
        # 有两个链接：形状级 + 运行级（同一个 url，classify_link 分类为 video）
        shape_link = [l for l in links if l["text"] == "观看演示"]
        assert len(shape_link) >= 1
        assert shape_link[0]["url"] == "https://example.com/demo.mp4"
        assert shape_link[0]["link_type"] == "video"
        # 验证 video Asset
        assert any(a.asset_type == AssetType.video_link for a in result.assets)

    def test_image_shape_with_hyperlink_links(self):
        """验证图片形状带超链接时 structured_data.links 正确输出。"""
        png_bytes = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
        )
        image_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        try:
            image_file.write(png_bytes)
            image_file.close()
            prs = Presentation()
            slide = _blank_slide(prs)
            slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "图片链接"
            pic = slide.shapes.add_picture(image_file.name, Inches(0.5), Inches(1.0))
            pic.click_action.hyperlink.address = "https://example.com/report.pdf"
            result = _parse(self.parser, _pptx_bytes(prs))
        finally:
            os.unlink(image_file.name)

        image_el = next(el for el in result.elements if el.element_type == ElementType.paragraph)
        assert image_el.structured_data is not None
        links = image_el.structured_data["links"]
        assert len(links) == 1
        assert links[0]["url"] == "https://example.com/report.pdf"
        assert links[0]["link_type"] == "document"
        # filename 作为链接文字 fallback
        assert links[0]["text"] != ""

    def test_mixed_hyperlink_and_plain_text(self):
        """验证多超链接混合文本（部分 run 有超链接、部分无）的正确处理。"""
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "混合文本"
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.0))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = "无链接 "
        r2 = p.add_run()
        r2.text = "视频链接"
        r2.hyperlink.address = "https://example.com/demo.mp4"
        r3 = p.add_run()
        r3.text = " 中间文字 "
        r4 = p.add_run()
        r4.text = "文档链接"
        r4.hyperlink.address = "https://example.com/doc.pdf"

        result = _parse(self.parser, _pptx_bytes(prs))
        paragraph = next(
            el for el in result.elements
            if el.element_type == ElementType.paragraph and "无链接" in el.text
        )
        # 文本包含所有 run 的文字
        assert "视频链接" in paragraph.text
        assert "中间文字" in paragraph.text
        assert "文档链接" in paragraph.text
        # links 仅包含有超链接的 run
        assert paragraph.structured_data is not None
        links = paragraph.structured_data["links"]
        assert len(links) == 2
        link_texts = {l["text"] for l in links}
        assert link_texts == {"视频链接", "文档链接"}
        link_types = {l["link_type"] for l in links}
        assert link_types == {"video", "document"}

    def test_list_item_hyperlinks_create_assets_and_links(self):
        """验证列表项中的超链接同时创建 Asset，链接写入列表容器 structured_data.links。"""
        prs = Presentation()
        slide = _blank_slide(prs)
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "资源列表"
        body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(6), Inches(1.5))
        # 第一条：带超链接
        p1 = body.text_frame.paragraphs[0]
        r1 = p1.add_run()
        r1.text = "下载文档"
        r1.hyperlink.address = "https://example.com/manual.pdf"
        # 第二条：无超链接
        p2 = body.text_frame.add_paragraph()
        p2.text = "纯文本说明"
        p2.level = 1

        result = _parse(self.parser, _pptx_bytes(prs))
        list_el = next(el for el in result.elements if el.element_type == ElementType.list)
        list_children = [el for el in result.elements if el.parent_element_id == list_el.element_id]

        # 链接信息在列表容器上
        assert list_el.structured_data is not None
        links = list_el.structured_data["links"]
        assert len(links) == 1
        assert links[0]["url"] == "https://example.com/manual.pdf"
        assert links[0]["text"] == "下载文档"

        # 验证 Asset 被创建（修复 gap）
        assert any(a.original_uri == "https://example.com/manual.pdf" for a in result.assets)

        # 子元素不包含 structured_data
        for child in list_children:
            assert child.structured_data is None or not child.structured_data.get("links")

    def test_classify_link_categories(self):
        """验证 classify_link 对各类 URL 的分类正确性。"""
        from parsers.utils import classify_link

        # 图片
        assert classify_link("https://example.com/photo.png") == "image"
        assert classify_link("https://example.com/icon.svg") == "image"
        # 视频（后缀）
        assert classify_link("https://example.com/demo.mp4") == "video"
        assert classify_link("https://example.com/clip.webm") == "video"
        # 视频（域名）
        assert classify_link("https://www.youtube.com/watch?v=abc") == "video"
        assert classify_link("https://www.bilibili.com/video/BV1xx") == "video"
        # 音频
        assert classify_link("https://example.com/song.mp3") == "audio"
        assert classify_link("https://example.com/podcast.wav") == "audio"
        # 文档
        assert classify_link("https://example.com/report.pdf") == "document"
        assert classify_link("https://example.com/data.xlsx") == "document"
        # 普通 URL
        assert classify_link("https://example.com/page") == "url"
        assert classify_link("https://example.com/about") == "url"

    def test_guess_mime_replacement_behavior(self):
        """验证 guess_mime 替换 _guess_mime 后行为不变。"""
        from parsers.utils import guess_mime

        # 使用 _pptx_bytes 创建简单 PPTX，验证图片 Asset 的 MIME 类型
        png_bytes = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
        )
        image_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        try:
            image_file.write(png_bytes)
            image_file.close()
            prs = Presentation()
            slide = _blank_slide(prs)
            slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5)).text = "MIME测试"
            slide.shapes.add_picture(image_file.name, Inches(0.5), Inches(1.0))
            result = _parse(self.parser, _pptx_bytes(prs))
        finally:
            os.unlink(image_file.name)

        assert len(result.assets) == 1
        assert result.assets[0].mime_type == "image/png"

        # 验证 URL Asset 的 MIME 推断
        assert guess_mime("https://example.com/doc.pdf", AssetType.document_link) == "application/pdf"
        assert guess_mime("https://example.com/unknown.xyz", AssetType.image) == "image/*"
